from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

    # Formatting Title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
def add_bullet_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

# --- SLIDE 1: Title ---
add_title_slide(
    "AI Readiness & Digital Maturity Dashboard",
    "Data-Driven Insights for Enterprise Transformation\nExecutive Summary"
)

# --- SLIDE 2: The Challenge ---
add_bullet_slide(
    "The Challenge",
    [
        "Organizations struggle to accurately measure current digital maturity.",
        "Difficulty in assessing readiness to adopt AI technologies effectively.",
        "Lack of structured data to evaluate adoption barriers.",
        "Blind spots regarding current capabilities and departmental needs."
    ]
)

# --- SLIDE 3: The Solution ---
add_bullet_slide(
    "The Solution",
    [
        "An automated pipeline and interactive dashboard.",
        "Simulates, processes, and visualizes assessment data across organizational personas.",
        "Granular insights across departments and seniority levels.",
        "Provides a centralized, highly visual view of the organization's digital posture."
    ]
)

# --- SLIDE 4: Key Features & Capabilities ---
add_bullet_slide(
    "Key Features & Capabilities",
    [
        "Automated Data Generation: Persona-based systems for highly realistic data.",
        "Intelligent Scoring Engine: Calculates Maturity and AI Readiness scores via complex logic.",
        "Dynamic Filtering: Assess data by Region, Department, and Role.",
        "Interactive Dashboard: Real-time exploration of KPIs via an intuitive web interface."
    ]
)

# --- SLIDE 5: Visual Insights ---
add_bullet_slide(
    "Visual Insights & Analytics",
    [
        "Maturity Distribution: Breakdown of capabilities across various axes.",
        "AI Adoption Barriers: Quantifying roadblocks like skill gaps, data security, and budget.",
        "Analytics Funnel: Tracking progression from initial capability to advanced AI deployment.",
        "Correlation Matrix: Identifying intersecting relationships between readiness factors."
    ]
)

# --- SLIDE 6: Business Impact ---
add_bullet_slide(
    "Business Impact",
    [
        "Enables precise data-driven decision-making for transformation strategies.",
        "Identifies 'low-hanging fruit' for immediate, high-ROI AI adoption.",
        "Highlights specific departments requiring targeted training and investment.",
        "Aligns technical initiatives with overarching business objectives."
    ]
)

# --- SLIDE 7: Conclusion & Next Steps ---
add_bullet_slide(
    "Conclusion & Next Steps",
    [
        "The dashboard is ready for deployment with real-world enterprise survey data.",
        "Next Steps:",
        "  1. Initiate Pilot Program with a select department.",
        "  2. Expand assessment criteria based on initial feedback.",
        "  3. Integrate live structured data sources for continuous monitoring."
    ]
)

# Save the presentation
prs.save("AI_Readiness_Executive_Summary.pptx")
print("Presentation generated successfully: AI_Readiness_Executive_Summary.pptx")
